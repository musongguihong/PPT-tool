from pptx import Presentation
from docx import Document

def extract_text_from_ppt(ppt_filename):
    prs = Presentation(ppt_filename)
    doc = Document()

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                if text:
                    doc.add_paragraph(text)

    return doc

def save_as_word(document, output_filename):
    document.save(output_filename)

def convert_ppt_to_word(ppt_filename):
    extracted_text = extract_text_from_ppt(ppt_filename)
    save_as_word(extracted_text, ppt_filename.split(".")[0] + ".docx")

# 使用方法：
convert_ppt_to_word("your_presentation.pptx")  # 替换成你的PPT文件名
