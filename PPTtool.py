from pptx import Presentation
from docx import Document

def extract_text_from_ppt(ppt_filename):
    prs = Presentation(ppt_filename)
    doc = Document()

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                if text:
                    doc.add_paragraph(text)

    return doc

def save_as_word(document, output_filename):
    document.save(output_filename)

def main():
    ppt_filenames = input("请拖入PPT文件并用空格分隔：").split()
    for ppt_filename in ppt_filenames:
        print(f"正在处理文件：{ppt_filename}")
        extracted_text = extract_text_from_ppt(ppt_filename)
        output_filename = ppt_filename.split(".")[0] + ".docx"
        save_as_word(extracted_text, output_filename)
        print(f"已将PPT文件中的文本保存为：{output_filename}")

    input("按 Enter 键退出...")

if __name__ == "__main__":
    main()
